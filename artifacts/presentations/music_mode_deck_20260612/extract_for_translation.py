from __future__ import annotations
"""Extract ordered paragraph records from the corrected EN deck + a numbered translatable list."""
import json, re
from pptx import Presentation

EN = "/tmp/MUSIC_MODE_INTRODUCTION_DECK_EN.pptx"
prs = Presentation(EN)

PAGE = re.compile(r'^\s*\d+\s*/\s*\d+\s*$')
def has_letter(s):
    return bool(re.search(r'[A-Za-z]', s))

records = []   # every paragraph: {i, s, sh, p, text, translate}
trans = []     # translatable subset: (i, text)
i = 0
for s_idx, slide in enumerate(prs.slides):
    for sh_idx, shape in enumerate(slide.shapes):
        if not shape.has_text_frame:
            continue
        for p_idx, para in enumerate(shape.text_frame.paragraphs):
            text = "".join(r.text for r in para.runs)
            if not text.strip():
                continue
            # auto-keep: page markers, and anything with no Latin letters (emoji, >, →, digits)
            translate = has_letter(text) and not PAGE.match(text)
            rec = {"i": i, "s": s_idx + 1, "sh": sh_idx, "p": p_idx, "text": text, "translate": translate}
            records.append(rec)
            if translate:
                trans.append((i, text))
            i += 1

with open("/tmp/en_records.json", "w") as f:
    json.dump(records, f, ensure_ascii=False, indent=0)

with open("/tmp/en_translatable.txt", "w") as f:
    for idx, text in trans:
        # one line per item; tabs/newlines flattened (none expected in a paragraph)
        f.write(f"{idx}\t{text}\n")

print(f"total paragraphs: {len(records)}  translatable: {len(trans)}")
