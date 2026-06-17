from __future__ import annotations
"""
apply_music_deck_corrections.py — committed SoT for the 2026-06-13 music-mode deck correction.

INPUT  : the operator's working copy (preserves every operator .pptx edit by construction).
OUTPUT : the corrected EN deck (operator edits + model corrections 1-8 layered on top).

Model corrections (see docs/specs/MUSIC_ONBOARDING_SONG_KIT_V1_SPEC.md §AMENDMENT-2026-06-13):
  1 both-ways data-driven mix (not a variant choice)   -> slides 2,3,4,7,12
  2 persona reuse vs composite (non-teacher path)        -> slides 2,6,15
  3 music mode = regular Pearl Prime ebook + music slots -> slides 2,3
  4 music_wrapper: 1P-default, music-attributed (#1527)  -> slides 3,8
  5 volume -> Pearl Prime ~800 + per-platform rollout    -> slides 10,11  (supersedes #1497 row-12 tiers)
  6 Behind-the-Song = podcast AND PDF                     -> slides 13,14
  7 no music-file distribution (music lives on channels)  -> slides 2,3,7,11,12,13,14
  8 slide-13 M3 YouTube/TikTok typos                      -> slide 13

Operator edits PRESERVED (not touched, except the M3 typo fix the operator asked for):
  s1 title "Pearl Prime / Music Content" + format list; s2 "on the Platform"; s3 "scripts"/no max-3;
  s4 "Ei Review"; s11 "Brand Director"/"Standard 800 book catalog"/"14 Markets"; s12 simplified
  Mode-Selector/intake; s13 M3 YouTube&TikTok + M5 "Podcasts:"; s15 "Composite USLF"/"spiritual wisdom".
"""
import sys
from pptx import Presentation


def _set_para(p, text):
    """Set a paragraph's text to `text`, preserving the first run's formatting."""
    runs = p.runs
    if not runs:
        r = p.add_run()
        r.text = text
        return
    runs[0].text = text
    for r in runs[1:]:
        r._r.getparent().remove(r._r)


def set_single(slide, shape_idx, text):
    """Replace a single-paragraph text frame with one paragraph `text` (keep run0 font)."""
    tf = slide.shapes[shape_idx].text_frame
    paras = tf.paragraphs
    _set_para(paras[0], text)
    for p in paras[1:]:
        p._p.getparent().remove(p._p)


def set_para(slide, shape_idx, contains, text):
    """In a multi-paragraph shape, replace the paragraph containing `contains`."""
    tf = slide.shapes[shape_idx].text_frame
    for p in tf.paragraphs:
        joined = "".join(r.text for r in p.runs)
        if contains in joined:
            _set_para(p, text)
            return
    raise SystemExit(f"  !! para match not found: shape {shape_idx} contains {contains!r}")


def apply(inp, outp):
    prs = Presentation(inp)
    S = prs.slides

    # ---- SLIDE 1 (idx 0): preserve operator title; complete the truncated subtitle ----
    set_single(S[0], 4, "How Pearl Prime Creates Books, Audiobooks, Podcasts & Manga at Global Scale")

    # ---- SLIDE 2 (idx 1): ebook+slots + data-mix + composite + no-music-file (corr 1,2,3,7) ----
    set_single(S[1], 8,
        "A regular Pearl Prime ebook + your music slots. Your data sets the mix — lyric books where "
        "you have lyrics, instrumental-reflection books where you don't, both in one brand. The "
        "teachings come from your own bank (if you teach) or the Pearl Prime composite. Your music "
        "stays on your channels; the book points to it.")
    set_single(S[1], 13,
        "Gen Z and Gen Alpha readers across 12 markets — US, Japan, Taiwan, China, and more — on "
        "KDP, Audible, Apple Books, and Google Play. Per-locale parity via Pearl_Localization. "
        "Addressed directly to the reader, in the author's voice.")

    # ---- SLIDE 3 (idx 2): music_wrapper note + Music-Format derived + companion-prompt reframe ----
    set_para(S[2], 7, "second-person framings",
        "(2P framings open & close) + a first-person, music-attributed wrapper on every music slot")
    set_single(S[2], 14,
        "Identity, themes, voice-craft, touchstones, healing intent, instruments, output "
        "preferences, Music-Format (data-derived)")
    set_single(S[2], 25, "Companion Prompt (optional)")
    set_single(S[2], 26,
        "An internal MusicGen mood-prompt per book for the brand-admin bundle — not a reader "
        "hand-out. Audio render is Phase-B-gated; your music is never distributed as a file.")

    # ---- SLIDE 4 (idx 3): fork->mix + drop max-3 + add voice lint (corr 1,4) ----
    set_single(S[3], 8,
        "The musician completes the Music Reflections Intake: themes, voice-craft, touchstones, "
        "healing intent, instruments, output preferences.")
    set_single(S[3], 22, "Mix: Lyrical + Instrumental")
    set_single(S[3], 23,
        "Your data sets the proportion — both coexist in one catalog. With-lyrics → lyric atoms "
        "(Pearl_Writer); no-lyrics → MusicGen mood-instruction text. No audio rendered in V1.")
    set_single(S[3], 32, "EI Review → Active")  # normalize operator's "Ei" to brand-standard "EI"
    set_single(S[3], 33,
        "Pearl_Editor reviews and promotes (Tier 1); the G1–G8 diversity gates + the "
        "music-attribution/voice lint run; injection unchanged.")

    # ---- SLIDE 6 (idx 5): add non-teacher composite path (corr 2) + grammar fix ----
    set_single(S[5], 13,
        "Ahjan already exists as a teacher, so the kit reuses that voice — no new persona. A "
        "non-teacher musician is welcome too: only the music is theirs; the teachings draw from "
        "the Pearl Prime composite.")
    set_para(S[5], 8, "reference future kits are measured",
        "The gold standard for every future musician kit.")  # QA: fix garbled dangling phrase

    # ---- SLIDE 7 (idx 6): both-coexist, data-driven mix; no companion-song (corr 1,7) ----
    set_single(S[6], 2, "Both Variants Coexist — Your Data Sets the Mix")
    set_single(S[6], 7,
        "Lyrical content at chapter opens, closes, and one mid-chapter beat. The more lyric "
        "material in your data, the larger this share of the catalog.")
    set_single(S[6], 13,
        "Reflection essays, journals, or meditations on your music at the same positions. A "
        "mostly-instrumental musician gets mostly these — plus a few lyric books — in one brand.")
    set_para(S[6], 14, "Companion song carries", "Your music — on your channels — carries the rest")

    # ---- SLIDE 8 (idx 7): music_wrapper rewrite, 1P+2P attributed examples (corr 4) ----
    set_single(S[7], 2, "The Author's Voice. Always Attributed to Your Music.")
    set_single(S[7], 3,
        "An author persona — pen-name or EI — references your music in its own voice, always "
        "attributed, never impersonating you. New music slots default to first person; existing "
        "second-person prose is kept.")
    set_single(S[7], 6, "EXAMPLE — FIRST PERSON (new-prose default)")
    set_single(S[7], 7,
        "“From Ahjan's music — on his guitar, a sustained chord, then a melody that lifts. "
        "That small lift of the spirit through sound is what I want to sit with you on, as we "
        "turn toward what weighs on you.”")
    set_single(S[7], 10, "Always attributed")
    set_single(S[7], 11,
        "Every music slot names your music and its source — “from Ahjan's music” — and "
        "states a true structure and a true benefit.")
    set_single(S[7], 14, "Second person — kept")
    set_single(S[7], 15,
        "“You put on Ahjan's track. The room gets small, then opens. The chord holds.” "
        "Where prose already reads in 2P, it is kept — still attributed.")
    # [18] "Refines over time" / [19] body — unchanged (on-model)

    # ---- SLIDE 10 (idx 9): kill volume tiers -> per-platform rollout (corr 5) ----
    set_single(S[9], 0, "QUALITY · PLATFORMS · ROLLOUT")
    set_single(S[9], 3,
        "Not AI slop. Deterministic books on your atom bank. Every brand targets ~800 — length "
        "routes to the platform it wins.")
    set_single(S[9], 12, "PER-PLATFORM ROLLOUT")
    set_single(S[9], 14, "FLAGSHIP")
    set_single(S[9], 15, "T7")
    set_single(S[9], 16,
        "~52k words. Wins full KDP ebook + Audible + Spotify audiobook at once — the only tier "
        "that wins the paid-book surfaces.")
    set_single(S[9], 18, "FUNNEL")
    set_single(S[9], 19, "T1–T5")
    set_single(S[9], 20,
        "Short Reads, podcast seasons, free AI-narration samples, social. The on-ramp tiers — "
        "mis-sold as books, ideal as funnel.")
    set_single(S[9], 22, "SEQUENCE")
    set_single(S[9], 23, "→")
    set_single(S[9], 24,
        "Flagship first (KDP + Audible), then funnel tiers fan across podcast / short-read / "
        "social. Volume follows Pearl Prime: ~800 high-confidence configs.")

    # ---- SLIDE 11 (idx 10): align tier chip + music-on-channels (corr 5,7) ----
    set_single(S[10], 17, "~800 Book Catalog")
    # [18] "Standard 800 book catalog" — operator edit, kept
    set_single(S[10], 27, "Music Stays on Your Channels")
    set_single(S[10], 28,
        "We never redistribute your music as files; each book links readers to it. An internal "
        "MusicGen prompt is optional, brand-admin only.")

    # ---- SLIDE 12 (idx 11): Music-Format derived + no companion-song bundling (corr 1,7) ----
    set_single(S[11], 23,
        "Now the Music Reflections Intake: themes, instruments (lyre / crystal bowls / frequency), "
        "and a data-derived Music-Format signal.")
    set_single(S[11], 33,
        "Catalog generator produces music-mode books; each book links readers to your music on "
        "your channels; Brand Director reviews & corrects; system improves.")

    # ---- SLIDE 13 (idx 12): M1 no-file reframe + M3 typo fix + M5 podcast+PDF (corr 6,7,8) ----
    set_single(S[12], 0, "AUDIO ON-RAMPS")
    set_single(S[12], 3,
        "MUSIC-MODE-FREEBIE-FUNNEL-V1-02 — on-ramps that point listeners to your music on your "
        "channels and pull them into your book catalog. We don't hand out your music as files.")
    set_single(S[12], 8, "Your Music, On Your Channels")
    set_single(S[12], 9,
        "No files handed out. Each book points readers to your music where it lives — your "
        "streaming, your channels — behind an email gate for the book funnel.")
    set_single(S[12], 10,
        "\"Hear 'Morning Room' on Ahjan's channel — the track this chapter was written around.\"")
    set_single(S[12], 22, "YouTube & TikTok")
    set_single(S[12], 23,
        "Inbound book sales from YouTube and TikTok. Outbound links in each book for readers to "
        "see and hear your music.")
    set_single(S[12], 24, "\"Get the 'Quiet Hours' links — 4 tracks that are the spine of this book.\"")
    set_single(S[12], 37,
        "A podcast episode AND a companion PDF — the deeper view of the song: what a track means, "
        "why it was written, what broke, what held. The reader starts understanding an artist.")

    # ---- SLIDE 14 (idx 13): funnel reframe — no song bundle; podcast+PDF (corr 6,7) ----
    set_single(S[13], 19, "Free On-Ramp")
    set_single(S[13], 20,
        "A pointer to your music on your channels (M1) or your YouTube/TikTok (M3) — email "
        "captured. They chose to hear more.")
    set_single(S[13], 38,
        "Behind-the-song (M5) — a podcast episode and a companion PDF. The musician's voice, "
        "unfiltered. The reader understands the artist.")
    set_single(S[13], 47,
        "Book CTA arrives naturally. The book is the deliverable; your music stays on your "
        "channels. Familiar territory — they already know the sound. They buy it.")
    set_single(S[13], 48, "THE BOOK")
    set_single(S[13], 58,
        "Books are the catalog. The on-ramps point to your music on your channels. Listeners who "
        "discover through sound become the most loyal readers.")

    # ---- SLIDE 15 (idx 14): operator edits kept verbatim — no change ----

    prs.save(outp)
    print(f"  saved -> {outp}")


if __name__ == "__main__":
    inp = sys.argv[1] if len(sys.argv) > 1 else "/Users/ahjan/Downloads/MUSIC_MODE_INTRODUCTION_DECK 3.pptx"
    outp = sys.argv[2] if len(sys.argv) > 2 else "/tmp/MUSIC_MODE_INTRODUCTION_DECK_EN.pptx"
    apply(inp, outp)
