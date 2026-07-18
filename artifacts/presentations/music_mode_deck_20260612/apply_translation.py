from __future__ import annotations
"""Validate a translations TSV against en_records.json and build the localized deck."""
import json, sys
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Pt

def _set_font(run, latin, ea):
    """Set BOTH the Latin (a:latin) and East-Asian (a:ea) typefaces on a run.
    LibreOffice selects CJK glyphs from a:ea — setting only a:latin leaves CJK as tofu."""
    run.font.name = latin                      # sets <a:latin>
    rPr = run._r.get_or_add_rPr()
    eael = rPr.find(qn('a:ea'))
    if eael is None:
        eael = rPr.makeelement(qn('a:ea'), {})
        latinel = rPr.find(qn('a:latin'))
        if latinel is not None:
            latinel.addnext(eael)              # schema order: latin, then ea
        else:
            rPr.append(eael)
    eael.set('typeface', ea)

def _set_para(p, text, latin, ea):
    runs = p.runs
    if not runs:
        r = p.add_run()
        r.text = text
        _set_font(r, latin, ea)
        return
    runs[0].text = text
    _set_font(runs[0], latin, ea)
    for r in runs[1:]:
        r._r.getparent().remove(r._r)


def build_serif_oracle(path):
    """positions (s,sh,p) whose run0 font is a serif title font in the existing localized deck."""
    prs = Presentation(path)
    serif = set()
    for s_idx, slide in enumerate(prs.slides, 1):
        for sh_idx, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            for p_idx, para in enumerate(shape.text_frame.paragraphs):
                name = None
                for r in para.runs:
                    name = r.font.name
                    break
                if name and ("Cambria" in name or "Mincho" in name or "Songti" in name or "Times" in name):
                    serif.add((s_idx, sh_idx, p_idx))
    return serif

def load_tsv(path):
    m = {}
    with open(path) as f:
        for ln, line in enumerate(f, 1):
            line = line.rstrip("\n")
            if not line.strip():
                continue
            if "\t" not in line:
                raise SystemExit(f"  !! {path}:{ln} no tab: {line[:60]!r}")
            idx, txt = line.split("\t", 1)
            m[int(idx)] = txt
    return m

def main(records_path, tsv_path, en_deck, out_deck, label, oracle_deck, serif_ea, sans_ea):
    records = json.load(open(records_path))
    trans = load_tsv(tsv_path)
    serif = build_serif_oracle(oracle_deck)
    SERIF_LATIN, SANS_LATIN = "Cambria", "Calibri"
    need = [r["i"] for r in records if r["translate"]]
    missing = [i for i in need if i not in trans]
    empty = [i for i in need if i in trans and not trans[i].strip()]
    extra = [i for i in trans if i not in set(r["i"] for r in records)]
    print(f"  [{label}] need={len(need)} have={len(trans)} missing={len(missing)} empty={len(empty)} extra={len(extra)}")
    if missing:
        raise SystemExit(f"  !! missing idx: {missing[:20]}")
    if empty:
        raise SystemExit(f"  !! empty idx: {empty[:20]}")

    prs = Presentation(en_deck)
    by_loc = {(r["s"], r["sh"], r["p"]): r["i"] for r in records}
    applied = 0
    for s_idx, slide in enumerate(prs.slides, 1):
        for sh_idx, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            for p_idx, para in enumerate(shape.text_frame.paragraphs):
                key = (s_idx, sh_idx, p_idx)
                if key in by_loc:
                    i = by_loc[key]
                    if i in trans and any(r["i"] == i and r["translate"] for r in records):
                        # only overwrite translatable paragraphs
                        rec = next(r for r in records if r["i"] == i)
                        if rec["translate"]:
                            is_serif = (s_idx, sh_idx, p_idx) in serif
                            latin = SERIF_LATIN if is_serif else SANS_LATIN
                            ea = serif_ea if is_serif else sans_ea
                            _set_para(para, trans[i], latin, ea)
                            applied += 1
    # slide-15 cover title: CJK glyphs are taller/wider than the EN Cambria → the big title wraps
    # to 2 lines and overlaps the subtitle. Shrink to fit one line (QA fix).
    title15 = prs.slides[14].shapes[2].text_frame.paragraphs[0]
    for r in title15.runs:
        r.font.size = Pt(36)
    prs.save(out_deck)
    print(f"  [{label}] applied {applied} translations -> {out_deck}")

if __name__ == "__main__":
    label, tsv, out, oracle, serif_font, sans_font = sys.argv[1:7]
    main("/tmp/en_records.json", tsv, "/tmp/MUSIC_MODE_INTRODUCTION_DECK_EN.pptx", out, label,
         oracle, serif_font, sans_font)
