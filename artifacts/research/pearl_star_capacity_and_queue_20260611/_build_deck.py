#!/usr/bin/env python3
"""Build PEARL_STAR_JOB_QUEUE_V1_DECK.pptx — 20 slides, Pearl Prime tokens.

Run once to generate the deck. Not a permanent artifact (the .pptx is).
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# Pearl Prime tokens
DARK_BG = RGBColor(0x0E, 0x0A, 0x06)   # near-black warm brown
LIGHT_BG = RGBColor(0xFA, 0xF6, 0xF0)  # parchment off-white
ACCENT = RGBColor(0xD9, 0x77, 0x06)    # warm orange-gold
TEXT_DARK = RGBColor(0x0E, 0x0A, 0x06)
TEXT_LIGHT = RGBColor(0xFA, 0xF6, 0xF0)
MUTED_DARK = RGBColor(0x4A, 0x3F, 0x35)
MUTED_LIGHT = RGBColor(0xCC, 0xC0, 0xB0)
DIVIDER_DARK = RGBColor(0x6B, 0x5A, 0x4A)

# Fonts — python-pptx writes the names; PowerPoint substitutes if not installed
F_HEAD = "Cormorant Garamond"
F_BODY = "DM Sans"
F_MONO = "DM Mono"
# Fallbacks PowerPoint will use if the above aren't on the viewer's box:
# Cormorant Garamond → Garamond/Georgia
# DM Sans → Calibri/Arial
# DM Mono → Consolas/Courier New

# 16:9 widescreen — 13.333 × 7.5 inches
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

OUT = Path("/Users/ahjan/phoenix_omega/.claude/worktrees/pearl-research-pearl-star-capacity-and-queue-v1-20260611/artifacts/research/pearl_star_capacity_and_queue_20260611/PEARL_STAR_JOB_QUEUE_V1_DECK.pptx")

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]  # blank layout

def fill_background(slide, rgb):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb

def add_text_box(slide, *, left, top, width, height, text, font=F_BODY, size=14,
                 color=TEXT_DARK, bold=False, italic=False, align=PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.vertical_anchor = anchor

    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.italic = italic
    return tb

def add_rect(slide, *, left, top, width, height, rgb, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = rgb
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    return sh

def add_accent_strip(slide, *, top, height, x=Inches(0.6), w=Inches(0.12)):
    return add_rect(slide, left=x, top=top, width=w, height=height, rgb=ACCENT)

def add_footer(slide, text, *, color=MUTED_DARK):
    add_text_box(slide, left=Inches(0.6), top=Inches(7.0), width=Inches(12.1),
                 height=Inches(0.4), text=text, font=F_MONO, size=9, color=color,
                 align=PP_ALIGN.LEFT)

def slide_number(slide, n, of=20, on_dark=False):
    color = MUTED_LIGHT if on_dark else MUTED_DARK
    add_text_box(slide, left=Inches(12.1), top=Inches(7.0), width=Inches(1.1),
                 height=Inches(0.4), text=f"{n} / {of}", font=F_MONO, size=9,
                 color=color, align=PP_ALIGN.RIGHT)

def content_header(slide, eyebrow, title):
    """Standard content-slide header layout."""
    # Eyebrow
    add_text_box(slide, left=Inches(0.6), top=Inches(0.45), width=Inches(12.1),
                 height=Inches(0.35), text=eyebrow.upper(), font=F_MONO, size=10,
                 color=ACCENT, bold=True)
    # Title
    add_text_box(slide, left=Inches(0.6), top=Inches(0.8), width=Inches(12.1),
                 height=Inches(0.9), text=title, font=F_HEAD, size=34, color=TEXT_DARK,
                 bold=True, line_spacing=1.0)

# ──────────────────────────────────────────────────────────────────────────
# Slide 1 — TITLE
s = prs.slides.add_slide(BLANK)
fill_background(s, DARK_BG)
# Vertical accent bar on left
add_rect(s, left=Inches(0.6), top=Inches(2.8), width=Inches(0.18),
         height=Inches(2.1), rgb=ACCENT)
add_text_box(s, left=Inches(1.1), top=Inches(2.55), width=Inches(11.5),
             height=Inches(0.5), text="PEARL_RESEARCH · OPERATOR BRIEF · 2026-06-11",
             font=F_MONO, size=11, color=ACCENT, bold=True)
add_text_box(s, left=Inches(1.1), top=Inches(3.05), width=Inches(11.5),
             height=Inches(1.8), text="Pearl Star Job Queue V1",
             font=F_HEAD, size=64, color=TEXT_LIGHT, bold=True, line_spacing=1.0)
add_text_box(s, left=Inches(1.1), top=Inches(4.5), width=Inches(11.5),
             height=Inches(0.6), text="Persistent · GPU-aware · Stall-recovering",
             font=F_HEAD, size=24, color=TEXT_LIGHT, italic=True)
add_text_box(s, left=Inches(1.1), top=Inches(5.4), width=Inches(11.5),
             height=Inches(0.5), text="Cap entry: PEARL-STAR-JOB-QUEUE-V1-01 (PROPOSAL)",
             font=F_MONO, size=12, color=MUTED_LIGHT)
slide_number(s, 1, on_dark=True)

# ──────────────────────────────────────────────────────────────────────────
# Slide 2 — THE PROBLEM (pull quote)
s = prs.slides.add_slide(BLANK)
fill_background(s, LIGHT_BG)
content_header(s, "the problem", "Why we need a queue")
# Big quote
add_text_box(s, left=Inches(1.2), top=Inches(2.1), width=Inches(11.0), height=Inches(0.6),
             text="“", font=F_HEAD, size=120, color=ACCENT, line_spacing=0.7)
add_text_box(s, left=Inches(2.0), top=Inches(2.3), width=Inches(10.0), height=Inches(3.5),
             text=("Every time we run something there's a hiccup — image, TTS, LLM jobs stall, "
                   "the server hangs, work disappears. I want to queue 1,000 book covers + 1,000 "
                   "manga panels + a podcast + LLM batch — and if the server reboots, the queue "
                   "picks up where it left off."),
             font=F_HEAD, size=22, color=TEXT_DARK, italic=True, line_spacing=1.3)
add_text_box(s, left=Inches(2.0), top=Inches(6.2), width=Inches(10.0), height=Inches(0.4),
             text="— Operator, May 2026", font=F_MONO, size=11, color=MUTED_DARK)
slide_number(s, 2)

# ──────────────────────────────────────────────────────────────────────────
# Slide 3 — PEARL STAR HARDWARE
s = prs.slides.add_slide(BLANK)
fill_background(s, LIGHT_BG)
content_header(s, "phase 1 · hardware", "Pearl Star at a glance")

# 4 KPI tiles
tiles = [
    ("1×", "RTX 5070 Ti GPU", "16 GB VRAM ceiling"),
    ("16", "CPU threads", "AMD Ryzen 7 7700 · 8c/16t"),
    ("64 GB", "RAM", "57 GiB available at idle"),
    ("1.5 TB", "Free NVMe", "Single drive · 17 % used"),
]
tx = Inches(0.6); ty = Inches(2.05); tw = Inches(3.0); th = Inches(2.2); gap = Inches(0.13)
for i, (big, mid, sub) in enumerate(tiles):
    x = tx + (tw + gap) * i
    add_rect(s, left=x, top=ty, width=tw, height=th, rgb=RGBColor(0xEF, 0xE8, 0xDC))
    add_rect(s, left=x, top=ty, width=tw, height=Inches(0.08), rgb=ACCENT)
    add_text_box(s, left=x+Inches(0.2), top=ty+Inches(0.25), width=tw-Inches(0.4),
                 height=Inches(0.9), text=big, font=F_HEAD, size=42, color=TEXT_DARK,
                 bold=True, line_spacing=1.0)
    add_text_box(s, left=x+Inches(0.2), top=ty+Inches(1.2), width=tw-Inches(0.4),
                 height=Inches(0.5), text=mid, font=F_BODY, size=14, color=TEXT_DARK,
                 bold=True)
    add_text_box(s, left=x+Inches(0.2), top=ty+Inches(1.55), width=tw-Inches(0.4),
                 height=Inches(0.6), text=sub, font=F_BODY, size=11, color=MUTED_DARK,
                 line_spacing=1.25)

# Bottom callout
add_rect(s, left=Inches(0.6), top=Inches(4.7), width=Inches(12.13), height=Inches(1.4),
         rgb=DARK_BG)
add_text_box(s, left=Inches(0.9), top=Inches(4.9), width=Inches(11.5),
             height=Inches(0.5), text="THE CONSTRAINT", font=F_MONO, size=10,
             color=ACCENT, bold=True)
add_text_box(s, left=Inches(0.9), top=Inches(5.25), width=Inches(11.5),
             height=Inches(0.9), text="One GPU = serial-by-default for VRAM-heavy work. The queue's job is to honor that constraint while keeping CPU-friendly work parallel.",
             font=F_HEAD, size=18, color=TEXT_LIGHT, italic=True, line_spacing=1.25)

add_footer(s, "Ubuntu 24.04 LTS · 47-day uptime · Tailscale-reached at pearlstar.tail7fd910.ts.net")
slide_number(s, 3)

# ──────────────────────────────────────────────────────────────────────────
# Slide 4 — CURRENT SOFTWARE STACK
s = prs.slides.add_slide(BLANK)
fill_background(s, LIGHT_BG)
content_header(s, "phase 1 · software", "What's installed today")

rows = [
    ("ComfyUI 0.18.1", "Image generation server · port 8188", "Manual-start (not systemd)"),
    ("Ollama systemd", "LLM server · port 11434", "gemma3:27b (17 GB) + qwen2.5:14b (9 GB)"),
    ("CosyVoice2 0.5B", "TTS · port 9880", "INSTALLED — not currently running"),
    ("Image models", "173 GB on disk", "flux-dev-fp8 · flux-schnell-fp8 · Qwen-Image · Animagine XL 4.0"),
    ("Custom nodes", "ComfyUI extensions", "PuLID Flux + SDXL · HandFixer"),
]
ry = Inches(2.05); rh = Inches(0.72); rg = Inches(0.08)
for i, (a, b, c) in enumerate(rows):
    y = ry + (rh + rg) * i
    add_rect(s, left=Inches(0.6), top=y, width=Inches(0.08), height=rh, rgb=ACCENT)
    add_text_box(s, left=Inches(0.85), top=y+Inches(0.05), width=Inches(3.0),
                 height=Inches(0.4), text=a, font=F_BODY, size=15, color=TEXT_DARK,
                 bold=True)
    add_text_box(s, left=Inches(0.85), top=y+Inches(0.38), width=Inches(3.0),
                 height=Inches(0.4), text=b, font=F_BODY, size=11, color=MUTED_DARK)
    add_text_box(s, left=Inches(4.0), top=y+Inches(0.15), width=Inches(8.7),
                 height=Inches(0.5), text=c, font=F_MONO, size=12, color=TEXT_DARK)

# Bottom line
add_rect(s, left=Inches(0.6), top=Inches(6.2), width=Inches(12.13), height=Inches(0.55),
         rgb=ACCENT)
add_text_box(s, left=Inches(0.9), top=Inches(6.27), width=Inches(11.5),
             height=Inches(0.45), text="No queue · no broker · blank slate for Phase A install",
             font=F_BODY, size=14, color=TEXT_DARK, bold=True, anchor=MSO_ANCHOR.MIDDLE)
slide_number(s, 4)

# ──────────────────────────────────────────────────────────────────────────
# Slide 5 — 4 WORKLOAD CLASSES
s = prs.slides.add_slide(BLANK)
fill_background(s, LIGHT_BG)
content_header(s, "decomposition", "The 4 workload classes")

cards = [
    ("T2I", "Text-to-Image",
     "flux-schnell · flux-dev (H1=A) · Qwen-Image · Animagine XL 4.0",
     "6–240 s per image · VRAM-heavy · cap = 1"),
    ("TTS", "Text-to-Speech",
     "CosyVoice2 (CJK) · Piper (EN)",
     "5–15 s per segment · CPU-friendly · cap = 2–8"),
    ("LLM", "Large Language Model",
     "Ollama qwen2.5:14b · gemma3:27b",
     "4–25 s per 200 tokens · VRAM-heavy · cap = 1–2"),
    ("ORCH", "Pipeline orchestration",
     "Pearl_Prime book · Manga V2 · Audiobook · Podcast · Video",
     "Multi-stage composition · CPU-only · cap = 2"),
]
cx = Inches(0.6); cy = Inches(2.05); cw = Inches(6.05); ch = Inches(2.0); cgx = Inches(0.13); cgy = Inches(0.16)
for i, (tag, name, engines, sizing) in enumerate(cards):
    col = i % 2; row = i // 2
    x = cx + (cw + cgx) * col
    y = cy + (ch + cgy) * row
    add_rect(s, left=x, top=y, width=cw, height=ch, rgb=RGBColor(0xEF, 0xE8, 0xDC))
    # Tag badge
    add_rect(s, left=x+Inches(0.2), top=y+Inches(0.2), width=Inches(1.0),
             height=Inches(0.5), rgb=ACCENT)
    add_text_box(s, left=x+Inches(0.2), top=y+Inches(0.2), width=Inches(1.0),
                 height=Inches(0.5), text=tag, font=F_BODY, size=14,
                 color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(s, left=x+Inches(1.4), top=y+Inches(0.25), width=cw-Inches(1.6),
                 height=Inches(0.5), text=name, font=F_HEAD, size=20, color=TEXT_DARK,
                 bold=True)
    add_text_box(s, left=x+Inches(0.25), top=y+Inches(0.9), width=cw-Inches(0.5),
                 height=Inches(0.5), text=engines, font=F_BODY, size=12,
                 color=TEXT_DARK)
    add_text_box(s, left=x+Inches(0.25), top=y+Inches(1.3), width=cw-Inches(0.5),
                 height=Inches(0.65), text=sizing, font=F_MONO, size=11,
                 color=MUTED_DARK)

add_text_box(s, left=Inches(0.6), top=Inches(6.55), width=Inches(12.13),
             height=Inches(0.4), text="Every Phoenix Omega job decomposes into one of these four classes",
             font=F_HEAD, size=14, color=MUTED_DARK, italic=True)
slide_number(s, 5)

# ──────────────────────────────────────────────────────────────────────────
# Slide 6 — CONCURRENCY MATRIX
s = prs.slides.add_slide(BLANK)
fill_background(s, LIGHT_BG)
content_header(s, "phase 2 · concurrency matrix", "What actually runs together")

# Safe column
add_rect(s, left=Inches(0.6), top=Inches(2.05), width=Inches(6.05), height=Inches(4.4),
         rgb=RGBColor(0xE6, 0xF0, 0xE0))
add_rect(s, left=Inches(0.6), top=Inches(2.05), width=Inches(6.05), height=Inches(0.55),
         rgb=RGBColor(0x2D, 0x6A, 0x33))
add_text_box(s, left=Inches(0.6), top=Inches(2.05), width=Inches(6.05),
             height=Inches(0.55), text="SAFE — run together", font=F_BODY, size=15,
             color=TEXT_LIGHT, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
safe_lines = [
    "1× flux + 1× CosyVoice2",
    "1× flux-dev + 1× CosyVoice2",
    "1× Qwen-Image + 1× CosyVoice2  (slow)",
    "1× Ollama qwen14b + 1× TTS",
    "4× CosyVoice2 in parallel",
    "8× Piper in parallel",
]
for i, line in enumerate(safe_lines):
    add_text_box(s, left=Inches(0.9), top=Inches(2.8) + Inches(0.55) * i,
                 width=Inches(5.5), height=Inches(0.45), text="✓  " + line,
                 font=F_MONO, size=13, color=TEXT_DARK)

# Unsafe column
add_rect(s, left=Inches(6.78), top=Inches(2.05), width=Inches(6.05), height=Inches(4.4),
         rgb=RGBColor(0xF4, 0xE2, 0xDC))
add_rect(s, left=Inches(6.78), top=Inches(2.05), width=Inches(6.05), height=Inches(0.55),
         rgb=RGBColor(0xA0, 0x32, 0x2D))
add_text_box(s, left=Inches(6.78), top=Inches(2.05), width=Inches(6.05),
             height=Inches(0.55), text="UNSAFE — serialize via /free",
             font=F_BODY, size=15, color=TEXT_LIGHT, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
unsafe_lines = [
    "2× flux concurrent on one GPU",
    "1× flux + Ollama gemma3:27b",
    "1× Qwen-Image + any LLM",
    "1× flux-dev + 1× Ollama qwen14b",
    "Anything VRAM-heavy + Qwen-Image",
    "Any two t2i engines concurrent",
]
for i, line in enumerate(unsafe_lines):
    add_text_box(s, left=Inches(7.08), top=Inches(2.8) + Inches(0.55) * i,
                 width=Inches(5.5), height=Inches(0.45), text="✗  " + line,
                 font=F_MONO, size=13, color=TEXT_DARK)

add_text_box(s, left=Inches(0.6), top=Inches(6.6), width=Inches(12.13),
             height=Inches(0.4), text="VRAM is the bottleneck — not threads.  Queue serializes VRAM-heavy work; runs CPU-heavy work concurrently.",
             font=F_HEAD, size=14, color=MUTED_DARK, italic=True)
slide_number(s, 6)

# ──────────────────────────────────────────────────────────────────────────
# Slide 7 — JOB SIZING RULE (the most important)
s = prs.slides.add_slide(BLANK)
fill_background(s, DARK_BG)
add_text_box(s, left=Inches(0.6), top=Inches(0.45), width=Inches(12.1),
             height=Inches(0.35), text="THE SIZING RULE", font=F_MONO, size=10,
             color=ACCENT, bold=True)

# Massive headline
add_text_box(s, left=Inches(0.6), top=Inches(1.3), width=Inches(12.13),
             height=Inches(2.3), text="One job =\none GPU dispatch",
             font=F_HEAD, size=80, color=TEXT_LIGHT, bold=True, line_spacing=1.0)

# Examples
examples = [
    ("1,000 book covers",     "→",  "1,000 jobs   (NOT 1)"),
    ("30-min podcast",        "→",  "60 × 30-s TTS jobs   (NOT 1)"),
    ("1 Pearl_Prime book",    "→",  "~30–100 atom jobs   (NOT 1)"),
    ("12 manga panels",       "→",  "12 jobs per chapter   (NOT 1)"),
]
ex_y = Inches(4.1)
for i, (a, arrow, b) in enumerate(examples):
    y = ex_y + Inches(0.5) * i
    add_text_box(s, left=Inches(1.5), top=y, width=Inches(4.3), height=Inches(0.45),
                 text=a, font=F_MONO, size=15, color=TEXT_LIGHT, align=PP_ALIGN.RIGHT)
    add_text_box(s, left=Inches(5.85), top=y, width=Inches(0.6), height=Inches(0.45),
                 text=arrow, font=F_MONO, size=15, color=ACCENT, bold=True,
                 align=PP_ALIGN.CENTER)
    add_text_box(s, left=Inches(6.55), top=y, width=Inches(6.2), height=Inches(0.45),
                 text=b, font=F_MONO, size=15, color=TEXT_LIGHT, bold=True)

add_text_box(s, left=Inches(0.6), top=Inches(6.85), width=Inches(12.13),
             height=Inches(0.35), text="Never bundle.  Always inspect per-unit.",
             font=F_HEAD, size=16, color=ACCENT, italic=True, align=PP_ALIGN.CENTER)
slide_number(s, 7, on_dark=True)

# ──────────────────────────────────────────────────────────────────────────
# Slide 8 — QUEUE RESEARCH TOP 3
s = prs.slides.add_slide(BLANK)
fill_background(s, LIGHT_BG)
content_header(s, "phase 3 · queue research", "Top 3 candidates side-by-side")

candidates = [
    ("PROCRASTINATE", "★ RECOMMENDED", "91 / 125",
     ["Postgres-backed", "ACID durability", "MIT license", "Async-first Python",
      "One dependency", "52K jobs/hr capacity"]),
    ("DRAMATIQ", "FALLBACK", "84 / 125",
     ["Redis-backed (Valkey BSD)", "Acks-when-done by default", "LGPL-3 library",
      "Simple ops", "Mature retry + DLQ", "Lighter than Postgres"]),
    ("CELERY", "DEFER", "67 / 125",
     ["Ack-early DEFAULT = foot-gun", "Requires careful config", "BSD license",
      "Largest ecosystem", "Battle-tested", "Heavy for solo operator"]),
]
ccx = Inches(0.6); ccy = Inches(2.05); ccw = Inches(4.04); cch = Inches(4.6); ccg = Inches(0.08)
for i, (name, tag, score, bullets) in enumerate(candidates):
    x = ccx + (ccw + ccg) * i
    is_winner = i == 0
    bg = RGBColor(0xEF, 0xE2, 0xC8) if is_winner else RGBColor(0xEF, 0xE8, 0xDC)
    add_rect(s, left=x, top=ccy, width=ccw, height=cch, rgb=bg)
    if is_winner:
        add_rect(s, left=x, top=ccy, width=ccw, height=Inches(0.08), rgb=ACCENT)
        add_rect(s, left=x, top=ccy+cch-Inches(0.08), width=ccw,
                 height=Inches(0.08), rgb=ACCENT)
    # Name + tag
    add_text_box(s, left=x+Inches(0.2), top=ccy+Inches(0.2), width=ccw-Inches(0.4),
                 height=Inches(0.5), text=name, font=F_HEAD, size=22,
                 color=TEXT_DARK, bold=True)
    add_text_box(s, left=x+Inches(0.2), top=ccy+Inches(0.75), width=ccw-Inches(0.4),
                 height=Inches(0.35), text=tag, font=F_MONO, size=11,
                 color=ACCENT if is_winner else MUTED_DARK, bold=True)
    # Score
    add_text_box(s, left=x+Inches(0.2), top=ccy+Inches(1.15), width=ccw-Inches(0.4),
                 height=Inches(0.7), text=score, font=F_HEAD, size=32,
                 color=TEXT_DARK, bold=True)
    # Bullets
    for j, b in enumerate(bullets):
        add_text_box(s, left=x+Inches(0.25), top=ccy+Inches(2.0) + Inches(0.4) * j,
                     width=ccw-Inches(0.5), height=Inches(0.35), text="•  " + b,
                     font=F_BODY, size=11, color=TEXT_DARK)

slide_number(s, 8)

# ──────────────────────────────────────────────────────────────────────────
# Slide 9 — WHY PROCRASTINATE WINS
s = prs.slides.add_slide(BLANK)
fill_background(s, LIGHT_BG)
content_header(s, "primary recommendation", "Why Procrastinate wins for Pearl Star")

reasons = [
    ("ACID durability", "Survives kill -9, server crash, power loss — Postgres WAL replay is the spec's hard requirement"),
    ("One dependency", "Postgres earns its keep beyond the queue: Pearl_PM tracker, dashboard, Pearl News metadata"),
    ("Python-async native", "Fits Pearl_Prime + audiobook comparator + manga dispatcher patterns out-of-box"),
    ("Commercial-clean license", "MIT — no contagion risk, no SSPL/AGPL ambiguity"),
    ("Proven scale headroom", "52,000 jobs/hr on single instance · Pearl Star's ceiling is ~10K jobs/day"),
]
ry = Inches(2.05); rh = Inches(0.8); rg = Inches(0.12)
for i, (head, sub) in enumerate(reasons):
    y = ry + (rh + rg) * i
    # number circle
    add_rect(s, left=Inches(0.6), top=y, width=Inches(0.6), height=Inches(0.6),
             rgb=ACCENT)
    add_text_box(s, left=Inches(0.6), top=y, width=Inches(0.6), height=Inches(0.6),
                 text=str(i+1), font=F_HEAD, size=22, color=TEXT_DARK, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(s, left=Inches(1.4), top=y, width=Inches(11.3), height=Inches(0.4),
                 text=head, font=F_BODY, size=15, color=TEXT_DARK, bold=True)
    add_text_box(s, left=Inches(1.4), top=y+Inches(0.35), width=Inches(11.3),
                 height=Inches(0.5), text=sub, font=F_BODY, size=12, color=MUTED_DARK,
                 line_spacing=1.25)

add_footer(s, "Sources: github.com/procrastinate-org/procrastinate · techplained.com/postgres-as-queue · 45 total citations in QUEUE_RESEARCH.md")
slide_number(s, 9)

# ──────────────────────────────────────────────────────────────────────────
# Slide 10 — GPU OVERLAY
s = prs.slides.add_slide(BLANK)
fill_background(s, LIGHT_BG)
content_header(s, "gpu overlay", "ComfyUI-Persistent-Queue (the belt-and-suspenders)")

# Flow boxes left to right
boxes = [
    ("PEARL_STAR_QUEUE", "Procrastinate + Postgres", "Durable work-list", LIGHT_BG, ACCENT),
    ("DISPATCH", "1-at-a-time", "/prompt + /free", LIGHT_BG, DIVIDER_DARK),
    ("COMFYUI + PQ", "Custom node", "Now persistent across restart", LIGHT_BG, ACCENT),
    ("RTX 5070 Ti", "16 GB VRAM", "The contended resource", DARK_BG, ACCENT),
]
bw = Inches(2.85); bh = Inches(2.3); by = Inches(2.6); bx = Inches(0.6); gap = Inches(0.27)
for i, (head, mid, sub, bg, border) in enumerate(boxes):
    x = bx + (bw + gap) * i
    add_rect(s, left=x, top=by, width=bw, height=bh, rgb=bg, line=border)
    tc = TEXT_LIGHT if bg == DARK_BG else TEXT_DARK
    mc = MUTED_LIGHT if bg == DARK_BG else MUTED_DARK
    add_rect(s, left=x, top=by, width=bw, height=Inches(0.08), rgb=ACCENT)
    add_text_box(s, left=x+Inches(0.15), top=by+Inches(0.3), width=bw-Inches(0.3),
                 height=Inches(0.5), text=head, font=F_MONO, size=11, color=ACCENT,
                 bold=True)
    add_text_box(s, left=x+Inches(0.15), top=by+Inches(0.85), width=bw-Inches(0.3),
                 height=Inches(0.6), text=mid, font=F_HEAD, size=18, color=tc, bold=True)
    add_text_box(s, left=x+Inches(0.15), top=by+Inches(1.6), width=bw-Inches(0.3),
                 height=Inches(0.6), text=sub, font=F_BODY, size=12, color=mc,
                 line_spacing=1.25)
    # Arrow between boxes
    if i < 3:
        ax = x + bw + Inches(0.02)
        add_text_box(s, left=ax, top=by+Inches(1.05), width=Inches(0.23),
                     height=Inches(0.5), text="▶", font=F_BODY, size=18, color=ACCENT,
                     bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_text_box(s, left=Inches(0.6), top=Inches(5.6), width=Inches(12.13),
             height=Inches(1.2),
             text=("ComfyUI's in-process queue is fast but loses jobs on Pearl Star reboot.\n"
                   "The community ComfyUI-Persistent-Queue custom node adds reboot-resume to it.\n"
                   "Pearl_Star_queue is the system-wide work-list; ComfyUI's queue is the GPU-arbiter slice."),
             font=F_BODY, size=14, color=TEXT_DARK, line_spacing=1.5)
slide_number(s, 10)

# ──────────────────────────────────────────────────────────────────────────
# Slide 11 — STALL DETECTION FLOW
s = prs.slides.add_slide(BLANK)
fill_background(s, LIGHT_BG)
content_header(s, "phase 4 · stall detection", "Heartbeat + Watchdog")

# Vertical flow boxes
steps = [
    ("00 : 00", "Worker emits heartbeat", "Every 30 s · phase, elapsed, VRAM, GPU util", RGBColor(0xEF, 0xE8, 0xDC)),
    ("01 : 00", "Watchdog tick", "Every 60 s · reads all heartbeat journals", RGBColor(0xEF, 0xE8, 0xDC)),
    ("02 : 00", "STALL_WARN logged", "elapsed > 2× expected · operator alert file-drop", RGBColor(0xF7, 0xE7, 0xC4)),
    ("03 : 00", "STALL_KILL fired", "elapsed > 3× expected · SIGTERM → 10 s → SIGKILL", RGBColor(0xEF, 0xCB, 0xB4)),
    ("∞", "CRASHED detected", "Heartbeat silent > 90 s · fresh worker dispatches", RGBColor(0xEA, 0xB7, 0xA0)),
]
sy = Inches(2.05); sh = Inches(0.78); sg = Inches(0.08); sx = Inches(0.6); sw = Inches(12.13)
for i, (clock, name, sub, bg) in enumerate(steps):
    y = sy + (sh + sg) * i
    add_rect(s, left=sx, top=y, width=sw, height=sh, rgb=bg)
    add_rect(s, left=sx, top=y, width=Inches(0.08), height=sh, rgb=ACCENT)
    # Clock chip
    add_text_box(s, left=sx+Inches(0.25), top=y+Inches(0.13), width=Inches(1.5),
                 height=Inches(0.5), text=clock, font=F_MONO, size=14, color=ACCENT,
                 bold=True)
    add_text_box(s, left=sx+Inches(1.9), top=y+Inches(0.08), width=Inches(4.5),
                 height=Inches(0.4), text=name, font=F_BODY, size=15, color=TEXT_DARK,
                 bold=True)
    add_text_box(s, left=sx+Inches(1.9), top=y+Inches(0.4), width=Inches(10.0),
                 height=Inches(0.4), text=sub, font=F_BODY, size=12, color=MUTED_DARK)

add_text_box(s, left=Inches(0.6), top=Inches(6.65), width=Inches(12.13),
             height=Inches(0.4),
             text="The operator's “every minute we check” becomes an automated SLA — no human polling needed.",
             font=F_HEAD, size=14, color=ACCENT, italic=True, align=PP_ALIGN.CENTER)
slide_number(s, 11)

# ──────────────────────────────────────────────────────────────────────────
# Slide 12 — AUTO-RECOVERY FLOW
s = prs.slides.add_slide(BLANK)
fill_background(s, LIGHT_BG)
content_header(s, "phase 4 · recovery", "Auto-recovery flow")

steps = [
    ("1", "Detect", "Watchdog observes stall or crash via heartbeat journal"),
    ("2", "Kill", "SIGTERM → 10 s grace → SIGKILL · verify exit"),
    ("3", "Verify VRAM reclaim", "Poll nvidia-smi until worker's footprint frees"),
    ("4", "Retry", "Mark auto_killed · apply retry policy per workload class"),
    ("5", "Dead-letter", "Retries exhausted → operator-review queue with full forensics"),
]
ry = Inches(2.0); rh = Inches(0.78); rg = Inches(0.12)
for i, (n, head, sub) in enumerate(steps):
    y = ry + (rh + rg) * i
    add_rect(s, left=Inches(0.6), top=y, width=Inches(0.78), height=rh, rgb=DARK_BG)
    add_text_box(s, left=Inches(0.6), top=y, width=Inches(0.78), height=rh,
                 text=n, font=F_HEAD, size=32, color=ACCENT, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(s, left=Inches(1.6), top=y+Inches(0.08), width=Inches(11.0),
                 height=Inches(0.4), text=head, font=F_BODY, size=16, color=TEXT_DARK,
                 bold=True)
    add_text_box(s, left=Inches(1.6), top=y+Inches(0.4), width=Inches(11.0),
                 height=Inches(0.4), text=sub, font=F_BODY, size=12, color=MUTED_DARK)

add_rect(s, left=Inches(0.6), top=Inches(6.55), width=Inches(12.13), height=Inches(0.5),
         rgb=ACCENT)
add_text_box(s, left=Inches(0.7), top=Inches(6.6), width=Inches(11.9),
             height=Inches(0.4),
             text="Retry budget per workload class:  t2i = 1  ·  llm = 2  ·  tts = 1  ·  orch = 3",
             font=F_MONO, size=12, color=TEXT_DARK, bold=True, anchor=MSO_ANCHOR.MIDDLE)
slide_number(s, 12)

# ──────────────────────────────────────────────────────────────────────────
# Slide 13 — pscli operator CLI
s = prs.slides.add_slide(BLANK)
fill_background(s, LIGHT_BG)
content_header(s, "operator control", "pscli — your queue command line")

# Terminal block
add_rect(s, left=Inches(0.6), top=Inches(2.0), width=Inches(12.13), height=Inches(4.4),
         rgb=DARK_BG)
add_rect(s, left=Inches(0.6), top=Inches(2.0), width=Inches(12.13), height=Inches(0.35),
         rgb=RGBColor(0x1F, 0x18, 0x10))
add_text_box(s, left=Inches(0.85), top=Inches(2.04), width=Inches(11.7),
             height=Inches(0.3), text="pearl_star · zsh",
             font=F_MONO, size=10, color=MUTED_LIGHT)
cmds = [
    "pscli status                       # queue depth + active workers + recent stalls",
    "pscli list --workload t2i --status pending",
    "pscli inspect <job_id>             # full job + heartbeat history + retries",
    "pscli pause                        # halt dispatch (current jobs complete)",
    "pscli resume                       # resume dispatch",
    "pscli drain                        # finish current; accept no new (pre-reboot)",
    "pscli kill <job_id>                # manual SIGTERM → SIGKILL",
    "pscli requeue <job_id>             # auto_killed / dead_letter → pending",
    "pscli vram-snapshot                # nvidia-smi diff vs last call",
    "pscli unload-comfyui               # POST /free to ComfyUI",
]
for i, c in enumerate(cmds):
    add_text_box(s, left=Inches(0.85), top=Inches(2.55) + Inches(0.36) * i,
                 width=Inches(11.5), height=Inches(0.35),
                 text="$ " + c, font=F_MONO, size=12, color=TEXT_LIGHT)

add_text_box(s, left=Inches(0.6), top=Inches(6.55), width=Inches(12.13),
             height=Inches(0.4),
             text="CLI from Day 1 · Web UI (Prometheus + Grafana) lands in Phase C",
             font=F_HEAD, size=14, color=MUTED_DARK, italic=True, align=PP_ALIGN.CENTER)
slide_number(s, 13)

# ──────────────────────────────────────────────────────────────────────────
# Slide 14 — PHASED ROLLOUT
s = prs.slides.add_slide(BLANK)
fill_background(s, LIGHT_BG)
content_header(s, "phase 6 · rollout", "Phased install — A → B → C → D")

phases = [
    ("A", "Foundation",
     "3-4 h · Pearl_Int session",
     ["Install Postgres 17 + Procrastinate",
      "Install ComfyUI-Persistent-Queue",
      "Add 6 systemd units (queue + watchdog + ComfyUI + CosyVoice2)",
      "Dogfood 1 workload (flux-schnell book covers)",
      "Acceptance: reboot-resume + stall + crash tests"]),
    ("B", "Four workloads",
     "4-6 h · Pearl_Dev session",
     ["Add llm + tts + orch worker pools",
      "Wire config/pipeline_registry.yaml stages",
      "Migrate manga / audiobook / podcast dispatchers",
      "Optional: vLLM as Ollama replacement",
      "Acceptance: composite scenario completes"]),
    ("C", "Operator dashboard",
     "6-8 h · Pearl_Dev + Pearl_Int",
     ["Install Prometheus + Grafana (self-hosted)",
      "Add nvidia_dcgm + node + queue exporters",
      "Provision Grafana dashboards",
      "Wire Pearl_PM tracker integration",
      "Optional: Prefect 3 web UI"]),
    ("D", "Multi-node",
     "FUTURE · only when Pearl Star 2 exists",
     ["Add Pearl Star 2 / 3 (additional GPU boxes)",
      "Multi-machine Procrastinate worker pool",
      "Optional: Ray for true GPU-aware scheduling",
      "Per-machine nvidia-smi reporting",
      "Not in V1 scope"]),
]
px = Inches(0.6); py = Inches(2.05); pw = Inches(3.0); ph = Inches(4.6); pg = Inches(0.17)
for i, (letter, name, when, items) in enumerate(phases):
    x = px + (pw + pg) * i
    bg = RGBColor(0xEF, 0xE2, 0xC8) if i < 3 else RGBColor(0xE8, 0xE3, 0xDC)
    add_rect(s, left=x, top=py, width=pw, height=ph, rgb=bg)
    # Big letter
    add_rect(s, left=x, top=py, width=pw, height=Inches(0.95), rgb=DARK_BG)
    add_text_box(s, left=x, top=py+Inches(0.1), width=pw, height=Inches(0.85),
                 text=letter, font=F_HEAD, size=44, color=ACCENT, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Name + when
    add_text_box(s, left=x+Inches(0.2), top=py+Inches(1.1), width=pw-Inches(0.4),
                 height=Inches(0.45), text=name, font=F_HEAD, size=18,
                 color=TEXT_DARK, bold=True)
    add_text_box(s, left=x+Inches(0.2), top=py+Inches(1.55), width=pw-Inches(0.4),
                 height=Inches(0.35), text=when, font=F_MONO, size=10,
                 color=ACCENT if i < 3 else MUTED_DARK, bold=True)
    # Items
    for j, it in enumerate(items):
        add_text_box(s, left=x+Inches(0.2), top=py+Inches(2.0) + Inches(0.5) * j,
                     width=pw-Inches(0.4), height=Inches(0.5), text="•  " + it,
                     font=F_BODY, size=10, color=TEXT_DARK, line_spacing=1.25)

slide_number(s, 14)

# ──────────────────────────────────────────────────────────────────────────
# Slide 15 — WHAT THIS UNBLOCKS
s = prs.slides.add_slide(BLANK)
fill_background(s, DARK_BG)
add_text_box(s, left=Inches(0.6), top=Inches(0.45), width=Inches(12.1),
             height=Inches(0.35), text="WHAT THIS UNBLOCKS", font=F_MONO, size=10,
             color=ACCENT, bold=True)

# Big headline
add_text_box(s, left=Inches(0.6), top=Inches(1.2), width=Inches(12.13),
             height=Inches(2.4),
             text="Queue 1,000 covers + 1,000 panels + a podcast + LLM batch.\n"
                  "Reboot the server mid-batch.\n"
                  "Every job completes.",
             font=F_HEAD, size=30, color=TEXT_LIGHT, bold=True, line_spacing=1.25)

# Use cases
uc = [
    ("Pearl_Prime $-maker scenario", "800 high-confidence configs all queued, dispatchable in any order"),
    ("Pearl News daily",            "25–75 LLM atoms queued · runs unattended overnight"),
    ("Audiobook + podcast",         "TTS-heavy work parallelized to CPU concurrency"),
    ("Manga V2 pipeline",           "Per-panel inspect-on-failure instead of chapter-level loss"),
]
ucy = Inches(4.5); uch = Inches(0.55); ucg = Inches(0.1)
for i, (head, sub) in enumerate(uc):
    y = ucy + (uch + ucg) * i
    add_rect(s, left=Inches(0.6), top=y, width=Inches(0.08), height=uch, rgb=ACCENT)
    add_text_box(s, left=Inches(0.85), top=y, width=Inches(4.5),
                 height=uch, text=head, font=F_BODY, size=14, color=TEXT_LIGHT,
                 bold=True, anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(s, left=Inches(5.5), top=y, width=Inches(7.3),
                 height=uch, text=sub, font=F_BODY, size=12, color=MUTED_LIGHT,
                 anchor=MSO_ANCHOR.MIDDLE)

slide_number(s, 15, on_dark=True)

# ──────────────────────────────────────────────────────────────────────────
# Slide 16 — OPEN OPERATOR QUESTIONS
s = prs.slides.add_slide(BLANK)
fill_background(s, LIGHT_BG)
content_header(s, "operator decisions", "Open questions — recommend defaults")

qs = [
    ("Q-PSQ-PRIMARY-QUEUE-01", "Primary queue framework",
     "Procrastinate + Postgres",
     "Alts: Dramatiq+Redis · Huey · ARQ"),
    ("Q-PSQ-BROKER-01", "Broker backend",
     "Postgres 17",
     "Alts: Redis (Valkey BSD) · NATS JetStream"),
    ("Q-PSQ-ROLLOUT-PHASE-A-WORKLOAD-01", "First dogfood workload",
     "flux-schnell (book covers)",
     "Alts: Qwen-Image · CosyVoice2 · Ollama"),
    ("Q-PSQ-DASHBOARD-01", "Operator dashboard phasing",
     "CLI Phase A → Web Phase C",
     "Alts: Web from Phase A · defer entirely"),
]
qy = Inches(2.05); qh = Inches(1.0); qg = Inches(0.12)
for i, (qid, q, default, alts) in enumerate(qs):
    y = qy + (qh + qg) * i
    add_rect(s, left=Inches(0.6), top=y, width=Inches(12.13), height=qh,
             rgb=RGBColor(0xEF, 0xE8, 0xDC))
    add_rect(s, left=Inches(0.6), top=y, width=Inches(0.08), height=qh, rgb=ACCENT)
    add_text_box(s, left=Inches(0.85), top=y+Inches(0.1), width=Inches(4.0),
                 height=Inches(0.35), text=qid, font=F_MONO, size=10, color=ACCENT,
                 bold=True)
    add_text_box(s, left=Inches(0.85), top=y+Inches(0.4), width=Inches(6.0),
                 height=Inches(0.5), text=q, font=F_BODY, size=14, color=TEXT_DARK,
                 bold=True, line_spacing=1.2)
    add_text_box(s, left=Inches(7.0), top=y+Inches(0.1), width=Inches(5.7),
                 height=Inches(0.45), text="DEFAULT: " + default,
                 font=F_BODY, size=12, color=TEXT_DARK, bold=True)
    add_text_box(s, left=Inches(7.0), top=y+Inches(0.5), width=Inches(5.7),
                 height=Inches(0.4), text=alts, font=F_MONO, size=10,
                 color=MUTED_DARK)

add_text_box(s, left=Inches(0.6), top=Inches(6.75), width=Inches(12.13),
             height=Inches(0.4),
             text="12 more in §9 of the V1 spec — Pearl_Research recommends defaults; operator decides.",
             font=F_HEAD, size=14, color=MUTED_DARK, italic=True, align=PP_ALIGN.CENTER)
slide_number(s, 16)

# ──────────────────────────────────────────────────────────────────────────
# Slide 17 — CROSS-REFERENCES
s = prs.slides.add_slide(BLANK)
fill_background(s, LIGHT_BG)
content_header(s, "composition", "Subsystems that route through the queue")

# Two columns
add_text_box(s, left=Inches(0.6), top=Inches(2.05), width=Inches(6.0),
             height=Inches(0.4), text="PHOENIX OMEGA PIPELINES", font=F_MONO, size=11,
             color=ACCENT, bold=True)
left_items = [
    ("Pearl_Prime ebook", "scripts/run_pipeline.py + atom decomposition"),
    ("Manga V2 (H1=A)", "scripts/manga/queue_panel_renders.py + 4-engine routing"),
    ("Audiobook", "Qwen-Only comparator loop + asyncio 24-concurrent"),
    ("Podcast", "scripts/podcast/ — 30-s CosyVoice2 chunks"),
    ("Video", "14-stage scripts/video/ pipeline"),
    ("Pearl News", "Pearl Star qwen14b CJK + Groq/Together EN fallback"),
]
for i, (a, b) in enumerate(left_items):
    add_text_box(s, left=Inches(0.6), top=Inches(2.5) + Inches(0.7) * i,
                 width=Inches(5.9), height=Inches(0.3), text=a,
                 font=F_BODY, size=13, color=TEXT_DARK, bold=True)
    add_text_box(s, left=Inches(0.6), top=Inches(2.8) + Inches(0.7) * i,
                 width=Inches(5.9), height=Inches(0.4), text=b,
                 font=F_MONO, size=11, color=MUTED_DARK)

add_text_box(s, left=Inches(6.8), top=Inches(2.05), width=Inches(6.0),
             height=Inches(0.4), text="PEARL STAR SUBSTRATE",
             font=F_MONO, size=11, color=ACCENT, bold=True)
right_items = [
    ("ComfyUI 0.18.1", "port 8188 · t2i dispatch via /prompt + /free"),
    ("Ollama systemd", "port 11434 · llm dispatch (qwen14b / gemma27b)"),
    ("CosyVoice2", "port 9880 · tts dispatch (cross-lingual)"),
    ("Piper", "binary · tts dispatch (English; CPU-only)"),
    ("Tailscale", "100.92.68.74 · Pearl Star network reach"),
    ("Postgres 17 (NEW)", "Phase A install · the queue's broker"),
]
for i, (a, b) in enumerate(right_items):
    add_text_box(s, left=Inches(6.8), top=Inches(2.5) + Inches(0.7) * i,
                 width=Inches(5.9), height=Inches(0.3), text=a,
                 font=F_BODY, size=13, color=TEXT_DARK, bold=True)
    add_text_box(s, left=Inches(6.8), top=Inches(2.8) + Inches(0.7) * i,
                 width=Inches(5.9), height=Inches(0.4), text=b,
                 font=F_MONO, size=11, color=MUTED_DARK)

add_text_box(s, left=Inches(0.6), top=Inches(6.85), width=Inches(12.13),
             height=Inches(0.4),
             text="The queue is the substrate — every existing subsystem composes through it.",
             font=F_HEAD, size=14, color=MUTED_DARK, italic=True, align=PP_ALIGN.CENTER)
slide_number(s, 17)

# ──────────────────────────────────────────────────────────────────────────
# Slide 18 — CAP-ENTRY PROPOSAL
s = prs.slides.add_slide(BLANK)
fill_background(s, LIGHT_BG)
content_header(s, "ratification", "Cap entry · PEARL-STAR-JOB-QUEUE-V1-01")

# Status badge
add_rect(s, left=Inches(0.6), top=Inches(2.05), width=Inches(2.5), height=Inches(0.55),
         rgb=ACCENT)
add_text_box(s, left=Inches(0.6), top=Inches(2.05), width=Inches(2.5),
             height=Inches(0.55), text="STATUS: PROPOSAL", font=F_MONO, size=12,
             color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text_box(s, left=Inches(3.2), top=Inches(2.05), width=Inches(9.5),
             height=Inches(0.55),
             text="→ (Pearl_Architect ratifies)  →  ACTIVE",
             font=F_MONO, size=12, color=MUTED_DARK, anchor=MSO_ANCHOR.MIDDLE)

# Decision rows
rows = [
    ("Decision", "Procrastinate + Postgres 17 primary · Dramatiq + Redis fallback · ComfyUI-Persistent-Queue GPU overlay"),
    ("Scope-in", "Single-box Pearl Star · 4 workload classes · reboot-resume durability · stall detection · pscli"),
    ("Scope-out", "Multi-node (Phase D) · paid LLM APIs (banned per CLAUDE.md) · K8s deployment · self-rolled broker"),
    ("Cross-refs", "spec + 6 research artifacts + skill manga_render_path · INTEGRATION_CREDENTIALS_REGISTRY §0 · pipeline_registry"),
    ("Next-action", "Operator reviews + answers Q-PSQ-* → Pearl_Architect ratifies → Pearl_Int Phase A install ws"),
    ("Owners", "Pearl_Architect (ratify) · Pearl_Dev (impl) · Pearl_Int (install) · Pearl_PM (tracker)"),
]
ry = Inches(2.85); rh = Inches(0.62); rg = Inches(0.08)
for i, (head, sub) in enumerate(rows):
    y = ry + (rh + rg) * i
    add_rect(s, left=Inches(0.6), top=y, width=Inches(0.08), height=rh, rgb=ACCENT)
    add_text_box(s, left=Inches(0.85), top=y+Inches(0.05), width=Inches(2.5),
                 height=Inches(0.5), text=head.upper(), font=F_MONO, size=11,
                 color=ACCENT, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(s, left=Inches(3.4), top=y, width=Inches(9.3),
                 height=rh, text=sub, font=F_BODY, size=12, color=TEXT_DARK,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)

slide_number(s, 18)

# ──────────────────────────────────────────────────────────────────────────
# Slide 19 — ACCEPTANCE CRITERIA
s = prs.slides.add_slide(BLANK)
fill_background(s, LIGHT_BG)
content_header(s, "phase a acceptance", "What 'done' looks like")

ac = [
    ("Reboot-resume",  "100 / 100 mixed jobs complete after Pearl Star reboot mid-batch · zero permanently lost"),
    ("Stall detection", "Injected sleep 600 → watchdog observes 2× / 3× thresholds → SIGKILL · retry succeeds"),
    ("Crash detection", "kill -9 on a worker mid-job → heartbeat-silent > 90 s → fresh worker → job completes"),
    ("VRAM reclaim",   "Post-AUTO-KILL · nvidia-smi shows VRAM at baseline within 30 s"),
    ("Operator CLI",   "pscli status returns < 2 s · pause / resume halt + restart dispatch correctly"),
    ("Tier policy",    "audit_llm_callers.py clean over Phase A queue code (no paid LLM API references)"),
]
ay = Inches(2.05); ah = Inches(0.7); ag = Inches(0.12)
for i, (head, sub) in enumerate(ac):
    y = ay + (ah + ag) * i
    add_rect(s, left=Inches(0.6), top=y, width=Inches(0.65), height=ah, rgb=DARK_BG)
    add_text_box(s, left=Inches(0.6), top=y, width=Inches(0.65), height=ah,
                 text=str(i+1), font=F_HEAD, size=24, color=ACCENT, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(s, left=Inches(1.45), top=y+Inches(0.05), width=Inches(11.3),
                 height=Inches(0.35), text=head, font=F_BODY, size=14,
                 color=TEXT_DARK, bold=True)
    add_text_box(s, left=Inches(1.45), top=y+Inches(0.35), width=Inches(11.3),
                 height=Inches(0.4), text=sub, font=F_MONO, size=11,
                 color=MUTED_DARK)

add_text_box(s, left=Inches(0.6), top=Inches(6.85), width=Inches(12.13),
             height=Inches(0.4),
             text="ALL 6 must pass before Phase A is marked DONE and Phase B begins.",
             font=F_HEAD, size=14, color=ACCENT, italic=True, align=PP_ALIGN.CENTER)
slide_number(s, 19)

# ──────────────────────────────────────────────────────────────────────────
# Slide 20 — DELIVERABLES + CLOSE
s = prs.slides.add_slide(BLANK)
fill_background(s, DARK_BG)
add_text_box(s, left=Inches(0.6), top=Inches(0.45), width=Inches(12.1),
             height=Inches(0.35), text="DELIVERABLES + NEXT", font=F_MONO, size=10,
             color=ACCENT, bold=True)

add_text_box(s, left=Inches(0.6), top=Inches(0.9), width=Inches(12.13),
             height=Inches(0.9), text="8 artifacts shipped in this PR",
             font=F_HEAD, size=32, color=TEXT_LIGHT, bold=True)

# Two-column delivery list
add_text_box(s, left=Inches(0.6), top=Inches(2.2), width=Inches(6.0),
             height=Inches(0.4), text="RESEARCH",
             font=F_MONO, size=11, color=ACCENT, bold=True)
left_files = [
    "HARDWARE_INVENTORY.md",
    "SOFTWARE_INVENTORY.md",
    "CONCURRENCY_MATRIX.md  (modeled)",
    "QUEUE_RESEARCH.md  (45 sources)",
    "STALL_RECOVERY_RUNBOOK.md",
    "JOB_SIZING_GUIDELINES.md",
]
for i, n in enumerate(left_files):
    add_text_box(s, left=Inches(0.6), top=Inches(2.6) + Inches(0.4) * i,
                 width=Inches(6.0), height=Inches(0.35), text="·  " + n,
                 font=F_MONO, size=12, color=TEXT_LIGHT)

add_text_box(s, left=Inches(6.8), top=Inches(2.2), width=Inches(6.0),
             height=Inches(0.4), text="SPEC + DECK",
             font=F_MONO, size=11, color=ACCENT, bold=True)
right_files = [
    "PEARL_STAR_JOB_QUEUE_V1_SPEC.md (16 §)",
    "PEARL_STAR_JOB_QUEUE_V1_DECK.pptx",
    "",
    "Cap entry: PEARL-STAR-JOB-QUEUE-V1-01",
    "  Status: PROPOSAL",
    "  Q-PSQ-* operator questions: 16",
]
for i, n in enumerate(right_files):
    add_text_box(s, left=Inches(6.8), top=Inches(2.6) + Inches(0.4) * i,
                 width=Inches(6.0), height=Inches(0.35),
                 text=("·  " + n) if n and not n.startswith("  ") else n,
                 font=F_MONO, size=12, color=TEXT_LIGHT)

# Closing line + accent rule
add_rect(s, left=Inches(0.6), top=Inches(5.5), width=Inches(2.0),
         height=Inches(0.04), rgb=ACCENT)
add_text_box(s, left=Inches(0.6), top=Inches(5.7), width=Inches(12.13),
             height=Inches(1.2),
             text=("Standing by for operator decisions on Q-PSQ-*\n"
                   "Phase A install kicks off after Pearl_Architect ratifies"),
             font=F_HEAD, size=22, color=TEXT_LIGHT, italic=True, line_spacing=1.4)
slide_number(s, 20, on_dark=True)

# ──────────────────────────────────────────────────────────────────────────
OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print(f"OK: deck saved at {OUT} ({OUT.stat().st_size:,} bytes; {len(prs.slides)} slides)")
