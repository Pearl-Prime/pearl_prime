#!/usr/bin/env python3
"""Build DURATION_PER_PLATFORM_DECK.pptx — Pearl Prime brand (near-black/gold/amber).
python-pptx only (Python 3.9 safe). Run: python3 build_deck.py
Token set + helpers reused from artifacts/qa/duration_marketing_targets_20260613/build_deck.py (house style)."""
import os
from pptx import Presentation
from pptx.util import Inches as I, Pt, Emu
from pptx.dml.color import RGBColor as C
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUT = os.path.dirname(os.path.abspath(__file__))
INK="1A1714"; INK2="2A2218"; CARD="332A1E"; GOLD="F5C842"; AMBER="F0A500"
WHITE="FFFFFF"; CITRUS="FF6B3D"; MIST="C5D1DB"; SLATE="7B8FA1"; GREEN="9CC79B"
EW, EH = I(13.333), I(7.5)
HDR="Georgia"; BODY="Calibri"

prs = Presentation(); prs.slide_width=EW; prs.slide_height=EH
BLANK = prs.slide_layouts[6]

def slide(bg=INK):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,EW,EH)
    r.fill.solid(); r.fill.fore_color.rgb=C.from_string(bg); r.line.fill.background()
    r.shadow.inherit=False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2,r._element)
    return s

def box(s,l,t,w,h,anchor=MSO_ANCHOR.TOP):
    tb=s.shapes.add_textbox(I(l),I(t),I(w),I(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=Pt(2);tf.margin_right=Pt(2);tf.margin_top=Pt(1);tf.margin_bottom=Pt(1)
    return tb,tf

def para(tf, text, size, color, bold=False, font=BODY, align=PP_ALIGN.LEFT,
         first=False, space_after=4, space_before=0, italic=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment=align; p.space_after=Pt(space_after); p.space_before=Pt(space_before)
    r=p.add_run(); r.text=text; f=r.font
    f.size=Pt(size); f.bold=bold; f.italic=italic; f.name=font
    f.color.rgb=C.from_string(color)
    return p

def runs(tf, parts, size, align=PP_ALIGN.LEFT, first=False, font=BODY, space_after=4):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment=align; p.space_after=Pt(space_after)
    for part in parts:
        if len(part)==4: text,sz,color,bold = part
        else: text,color,bold = part; sz=size
        r=p.add_run(); r.text=text; r.font.size=Pt(sz); r.font.bold=bold
        r.font.name=font; r.font.color.rgb=C.from_string(color)
    return p

def rect(s,l,t,w,h,fill,line=None,shape=MSO_SHAPE.ROUNDED_RECTANGLE,line_w=1.0):
    sp=s.shapes.add_shape(shape,I(l),I(t),I(w),I(h))
    sp.shadow.inherit=False
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=C.from_string(fill)
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=C.from_string(line); sp.line.width=Pt(line_w)
    return sp

def footer(s,page):
    _,tf=box(s,0.6,7.04,11,0.35)
    para(tf,"Pearl Research  ·  Duration-Per-Platform Plan  ·  2026-06-13  ·  plan for operator review (config PR do-NOT-merge)",
         9,MIST,first=True,space_after=0)
    _,tf2=box(s,12.3,7.04,0.6,0.35)
    para(tf2,str(page),11,GOLD,bold=True,align=PP_ALIGN.RIGHT,first=True,space_after=0)

def title(s,text,sub=None):
    _,tf=box(s,0.6,0.70,12.1,1.05)
    para(tf,text,32,GOLD,bold=True,font=HDR,first=True,space_after=2)
    if sub: para(tf,sub,14,MIST,italic=True)

def tag(s,text,col=AMBER):
    w=max(0.95,0.30+0.105*len(text))
    sp=rect(s,0.6,0.34,w,0.32,None,line=col,line_w=1.25)
    tf=sp.text_frame; tf.margin_top=0;tf.margin_bottom=0
    para(tf,text,11,col,bold=True,align=PP_ALIGN.CENTER,first=True,space_after=0)

def style_cell(cell,text,size,color,fill,bold=False,align=PP_ALIGN.LEFT,font=BODY):
    cell.fill.solid(); cell.fill.fore_color.rgb=C.from_string(fill)
    cell.vertical_anchor=MSO_ANCHOR.MIDDLE
    cell.margin_left=Pt(6);cell.margin_right=Pt(5);cell.margin_top=Pt(1);cell.margin_bottom=Pt(1)
    tf=cell.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=text; f=r.font
    f.size=Pt(size); f.bold=bold; f.name=font; f.color.rgb=C.from_string(color)

def add_table(s,l,t,w,col_w,rows,row_h,header_fill=AMBER,header_fg=INK,
              body_size=11,cell_colors=None,text_colors=None,aligns=None):
    nr=len(rows); nc=len(rows[0])
    gf=s.shapes.add_table(nr,nc,I(l),I(t),I(w),I(row_h*nr))
    tbl=gf.table
    tblPr=tbl._tbl.tblPr
    for a in ('firstRow','bandRow'):
        tblPr.set(a,'0')
    for ci,cw in enumerate(col_w): tbl.columns[ci].width=I(cw)
    for ri in range(nr):
        tbl.rows[ri].height=I(row_h)
        for cj in range(nc):
            cell=tbl.cell(ri,cj)
            if ri==0:
                style_cell(cell,rows[ri][cj],body_size+1,header_fg,header_fill,bold=True,
                           align=(aligns[cj] if aligns else PP_ALIGN.LEFT),font=HDR)
            else:
                fill = INK2 if ri%2 else CARD
                if cell_colors and (ri,cj) in cell_colors: fill=cell_colors[(ri,cj)]
                fg = WHITE
                if text_colors and (ri,cj) in text_colors: fg=text_colors[(ri,cj)]
                bold = bool(text_colors and (ri,cj) in text_colors)
                style_cell(cell,rows[ri][cj],body_size,fg,fill,bold=bold,
                           align=(aligns[cj] if aligns else PP_ALIGN.LEFT))
    return tbl

# ───────────────────────── Slide 1 — Title ─────────────────────────
s=slide(INK)
rect(s,0,0,0.32,7.5,GOLD,shape=MSO_SHAPE.RECTANGLE)
_,tf=box(s,1.1,2.0,11.5,3.1)
para(tf,"The Duration-Per-Platform Plan",44,GOLD,bold=True,font=HDR,first=True,space_after=6)
para(tf,"The length that WINS on each surface — mapped to our 7 honest tiers.",19,WHITE,space_after=2)
para(tf,"And the hard truth: at our REAL render lengths, we win none of the flagship paid-book surfaces yet.",19,WHITE,space_after=14)
runs(tf,[("Pearl Research      ",13,MIST,False),
         ("50+ cited sources · plan for operator review · config PR do-NOT-merge",13,AMBER,True)],13)
_,tf2=box(s,1.1,5.95,11,0.6)
para(tf2,"Builds on the merged honest ladder (PR #1550) + CDIS §4 · validates/refreshes both with current 2026 research",12,MIST,italic=True,first=True)

# ───────────────────────── Slide 2 — the gap ─────────────────────────
s=slide(INK); tag(s,"THE GAP",CITRUS); title(s,"We render a podcast episode; the money is in a 6-hour book")
rect(s,0.6,1.95,5.9,4.45,CARD,line=CITRUS,line_w=1.5)
_,tf=box(s,0.95,2.18,5.25,1.9)
para(tf,"what we RENDER today",13,MIST,first=True,space_after=2)
para(tf,"~33 min",52,CITRUS,bold=True,font=HDR,space_after=2)
para(tf,"thin catalog book ≈ 5,000 words ≈ 20 pp.",14,WHITE,space_after=2)
para(tf,"Gold/depth-fill path tops out ~2.4 hr (22k words).",13,MIST,space_after=0)
_,tf=box(s,0.95,4.35,5.25,1.9)
para(tf,"its ONLY winning home:",13,MIST,first=True,space_after=2)
para(tf,"one podcast episode",24,GREEN,bold=True,font=HDR,space_after=4)
para(tf,"On every PAID-book surface it is a lead-magnet or a “short read.”",13,WHITE,space_after=0)
rect(s,6.85,1.95,5.9,4.45,INK2,line=AMBER,line_w=1.25)
_,tf=box(s,7.2,2.18,5.25,4.0)
para(tf,"what WINS on Audible / KDP",13,MIST,first=True,space_after=2)
para(tf,"5–7 hr  ·  150–230 pp",40,GOLD,bold=True,font=HDR,space_after=6)
para(tf,"≈ 45,000–63,000 words. Only our top tier (T7 “Complete,” 52k) reaches it — and nothing renders that long today.",16,WHITE,space_after=12)
runs(tf,[("The labels are now honest (#1550).  ",14,GOLD,True),
         ("The content is too short to win the surfaces that pay the most.",14,WHITE,False)],14)
footer(s,2)

# ───────────────────────── Slide 3 — tier→platform fit matrix ─────────────────────────
s=slide(INK); tag(s,"THE CROSS-WALK"); title(s,"Which tier wins where — the missing map","Our 7 tiers are NOT interchangeable across surfaces. ✓✓ flagship win · ✓ wins · ~ wrong home · ✗ loser")
rows=[["Tier","words","audio","pp","KDP full","Audible","Spotify a/b","Podcast"],
 ["T1 Quick Reset","3.5k","23m","14","✗","✗","✗","✓✓ 1 ep"],
 ["T2 Mini","4.5k","30m","18","✗","✗","✗","✓✓ 1 ep"],
 ["T3 Short","6k","40m","24","~ short-read","✗","✗","✓✓ 1 ep"],
 ["T4 One-Hour","9k","60m","36","~ short-read","✗","✗","✓ special"],
 ["T5 Standard","22k","2.4h","88","~ short-read","~ <3h floor","~ low","✓✓ season"],
 ["T6 Long-Form","30k","3.3h","120","~ low-end","✓ 3–5h","~","✓ season"],
 ["T7 Complete","52k","5.8h","208","✓✓","✓✓","✓","✓ multi"],
]
tc={}; cc={}
# green the T7 flagship wins; citrus the losers column-ish
for cj in range(4,8):
    tc[(7,cj)]=GREEN
for ri,vals in [(1,[5,6,7]),(2,[5,6,7]),(3,[6,7]),(4,[6,7]),(5,[]),(6,[]),(7,[])]:
    for cj in vals: tc[(ri,cj)]=CITRUS
add_table(s,0.6,1.9,12.13,[2.0,0.95,0.95,0.7,1.9,1.75,1.75,2.13],rows,0.54,body_size=11,
          text_colors=tc,aligns=[PP_ALIGN.LEFT,PP_ALIGN.CENTER,PP_ALIGN.CENTER,PP_ALIGN.CENTER,PP_ALIGN.CENTER,PP_ALIGN.CENTER,PP_ALIGN.CENTER,PP_ALIGN.CENTER])
_,tf=box(s,0.6,6.42,12.1,0.5)
runs(tf,[("Only T7 wins full KDP + Audible + Spotify at once.  ",12.5,GOLD,True),
         ("T1–T5 are podcast / short-read / lead-magnet plays — structural losers as standalone audiobooks.",12.5,WHITE,False)],12.5,first=True)
footer(s,3)

# ───────────────────────── Slide 4 — per-platform targets ─────────────────────────
s=slide(INK); tag(s,"THE TARGETS"); title(s,"What each surface rewards — the number to hit + why")
rows=[["Platform","Target to WIN","Our tier","Why (mechanism)"],
 ["Amazon KDP (full)","150–230 pp / 35–50k w","T7","KU pays per page READ × rising rate; <150pp = short-read"],
 ["Audible / ACX","5–7 hr / 45–63k w","T7 (T6 entry)","length sets price tier + consumption-pool weight; <3h = bottom"],
 ["Spotify audiobook","4–6 hr / 36–54k w","T7","15 hr/mo at normal speed → tryable; payout on COMPLETED hrs"],
 ["Spotify/Apple podcast","20–40 min / episode","T1–T3 = 1 ep","completion flat 10→45 min; cadence ≈ +5 chart spots"],
 ["Apple/Google AI a/b","length-agnostic (EN only)","any (T7 best)","flat per-sale → length = price justification, not payout"],
 ["YouTube long / Shorts","8–15 min / 15–60s (max 180)","video","watch-time × satisfaction; 8 min = mid-roll; trailer 60–90s"],
 ["WEBTOON / Piccoma","~8–12 min read; 25–30 ep","manga","chapter count + cadence drive unlock revenue, not panels"],
]
add_table(s,0.6,1.9,12.13,[2.45,2.7,1.7,5.28],rows,0.56,body_size=11,
          aligns=[PP_ALIGN.LEFT,PP_ALIGN.LEFT,PP_ALIGN.CENTER,PP_ALIGN.LEFT])
_,tf=box(s,0.6,6.46,12.1,0.5)
para(tf,"Every read-through surface (KU, Audible-new, Spotify, Kobo) pays on COMPLETED length — win with the longest length we can fill at ≥70% completion. Never padding.",11,MIST,italic=True,first=True)
footer(s,4)

# ───────────────────────── Slide 5 — meets/misses ─────────────────────────
s=slide(INK); tag(s,"VERIFICATION"); title(s,"Our REAL durations vs the platform targets","✓ wins · ~ marginal/wrong-home · ✗ miss · ✗✗ catastrophic")
rows=[["Platform target","THIN ~5k/33m","GOLD ~22k/2.4h","Ladder T7 52k"],
 ["KDP full ebook (150–230pp)","✗✗ lead-magnet","✗ short-read (88pp)","✓ MEETS (208pp)"],
 ["Audible (5–7h; floor 3h)","✗✗ bottom tier","✗ under 3h floor","✓ MEETS (5.8h)"],
 ["Spotify audiobook (4–6h)","✗✗","✗ below 4h","✓ MEETS (5.8h)"],
 ["Podcast episode (20–40m)","✓ = 1 episode","✓ = a season","✓ multi-season"],
 ["Apple/Google AI (EN)","~ 33-min sample","✓ listable","✓ best value"],
]
tc={}
for cj in [1,2]:
    tc[(1,cj)]=CITRUS; tc[(2,cj)]=CITRUS; tc[(3,cj)]=CITRUS
tc[(4,1)]=GREEN; tc[(4,2)]=GREEN
for ri in range(1,4): tc[(ri,3)]=GREEN
tc[(4,3)]=GREEN
add_table(s,0.6,2.05,12.13,[3.5,2.85,2.85,2.93],rows,0.66,body_size=11.5,
          text_colors=tc,aligns=[PP_ALIGN.LEFT,PP_ALIGN.CENTER,PP_ALIGN.CENTER,PP_ALIGN.CENTER])
_,tf=box(s,0.6,6.65,12.1,0.5)
runs(tf,[("The thin render's ONLY ✓ is a podcast episode.  ",12.5,GREEN,True),
         ("Today we ship zero product that wins a full Audible audiobook or full KDP book — the T7 tier exists but no render reaches it.",12.5,WHITE,False)],12.5,first=True)
footer(s,5)

# ───────────────────────── Slide 6 — the config bug ─────────────────────────
s=slide(INK); tag(s,"ROOT CAUSE",CITRUS); title(s,"The routing was built on the OLD wrong label")
rect(s,0.6,1.95,5.95,4.5,CARD,line=CITRUS,line_w=1.5)
_,tf=box(s,0.95,2.18,5.3,4.1)
para(tf,"platform_knob_tuning.yaml — before",13,CITRUS,bold=True,first=True,space_after=6)
para(tf,"audible:",14,MIST,font=HDR,space_after=1)
para(tf,"  preferred_runtimes:",13,WHITE,space_after=1)
para(tf,"    [standard_book, extended_book_2h]",13,CITRUS,bold=True,space_after=1)
para(tf,"  ideal_runtime_hours: 5",13,WHITE,space_after=8)
para(tf,"…routes 2.3–2.5 hr books to Audible while STATING the ideal is 5 hr. Authored when standard_book was mislabeled “55 min.”",13.5,MIST,italic=True,space_after=0)
rect(s,6.9,1.95,5.85,4.5,INK2,line=GOLD,line_w=1.25)
_,tf=box(s,7.25,2.18,5.2,4.1)
para(tf,"after — Pearl Research correction",13,GOLD,bold=True,first=True,space_after=6)
para(tf,"audible:",14,MIST,font=HDR,space_after=1)
para(tf,"  preferred_runtimes:",13,WHITE,space_after=1)
para(tf,"    [deep_book_6h, deep_book_4h]",13,GREEN,bold=True,space_after=8)
para(tf,"Only deep_book_6h (5.8 hr) + deep_book_4h (3.3 hr) clear the Audible 3-hr floor and reach the stated ideal.",13.5,WHITE,space_after=8)
para(tf,"#1550 fixed the label (55→147); this config was never reconciled. Same fix on Spotify/Apple/Google/Kobo + findaway (now INaudio).",12.5,MIST,italic=True,space_after=0)
footer(s,6)

# ───────────────────────── Slide 7 — the fix split ─────────────────────────
s=slide(INK); tag(s,"THE FIX"); title(s,"Two fixes, cleanly separated")
rect(s,0.6,1.95,5.95,4.5,INK2,line=GOLD,line_w=1.25)
_,tf=box(s,0.95,2.2,5.3,4.1)
para(tf,"CONFIG  —  done in-lane (this PR)",16,GOLD,bold=True,font=HDR,first=True,space_after=8)
para(tf,"• Re-route Audible/Spotify/Apple/Google/Kobo to the tiers that reach their stated ideal (deep_book_6h/4h).",13.5,WHITE,space_after=6)
para(tf,"• Cite + correct the platform profiles (Shorts 60→180s, podcast 20→40 min, Findaway→INaudio, Audible consumption model).",13.5,WHITE,space_after=6)
para(tf,"• Add the tier→platform cross-walk as machine-readable SSOT.",13.5,WHITE,space_after=8)
para(tf,"RULE-0 clean: additive/corrective, 0 deletions. DO-NOT-MERGE pending operator review.",12.5,MIST,italic=True,space_after=0)
rect(s,6.9,1.95,5.85,4.5,CARD,line=AMBER,line_w=1.25)
_,tf=box(s,7.25,2.2,5.2,4.1)
para(tf,"CONTENT-LENGTH  —  routed (not faked)",16,AMBER,bold=True,font=HDR,first=True,space_after=8)
para(tf,"• To win Audible/KDP at all, we need a real 52k-word “Complete” (T7) flagship.",14,WHITE,space_after=6)
para(tf,"• Fed to the WRITING PROGRAM as a word_target — the registry already defines it (deep_book_6h).",13.5,WHITE,space_after=6)
para(tf,"• Close the thin-render shortfall (renders ~5k vs target) with real depth, ≥70% completion.",13.5,WHITE,space_after=8)
runs(tf,[("Never pad. ",14,CITRUS,True),
         ("Padding loses the read-through revenue that is the whole reason to be long.",13,WHITE,False)],13,space_after=0)
footer(s,7)

# ───────────────────────── Slide 8 — CJK provisional ─────────────────────────
s=slide(INK); tag(s,"CJK — PROVISIONAL"); title(s,"Char-based, not word-based — gate on a measured render")
rows=[["Locale","T7 audiobook","ebook vs local norm","Confidence"],
 ["en-US","5.8 hr","208pp in 150–230 ✓","baseline"],
 ["ja-JP","6.2 hr (×2.15 ÷300)","111.8k 文字 in 80–120k ✓","MEASURED + rate hard-sourced"],
 ["zh (CN/TW/SG)","7.3 hr (slow pace)","83k 字 BELOW 100–150k → wants MORE","1.6 UNVERIFIED — measure"],
 ["ko-KR","~5.0 hr","ship FULL + Millie summary","2.0 + rate UNVERIFIED — measure"],
]
tc={(2,3):GREEN,(3,3):CITRUS,(4,3):CITRUS,(2,0):GOLD,(3,0):GOLD,(4,0):GOLD}
add_table(s,0.6,2.0,12.13,[1.6,3.3,4.2,3.03],rows,0.7,body_size=11.5,
          text_colors=tc,aligns=[PP_ALIGN.LEFT,PP_ALIGN.CENTER,PP_ALIGN.LEFT,PP_ALIGN.LEFT])
_,tf=box(s,0.6,6.05,12.1,1.0)
runs(tf,[("ja narration 300 文字/分 is the one hard CJK anchor.  ",12.5,GOLD,True),
         ("zh wants a HIGHER word_target (Ximalaya = episodic). Korea = ship BOTH a Millie summary + a Welaaa full edition.",12.5,WHITE,False)],12.5,first=True,space_after=4)
para(tf,"Do NOT ship any CJK platform duration until one zh + one ko render is measured on Pearl Star (CosyVoice2 currently off).",12,CITRUS,italic=True)
footer(s,8)

# ───────────────────────── Slide 9 — next actions ─────────────────────────
s=slide(INK); tag(s,"DECISION")
_,tf=box(s,0.6,0.7,12.1,1.0)
para(tf,"What we need from you",32,GOLD,bold=True,font=HDR,first=True,space_after=2)
para(tf,"The labels are honest. Now decide the platform strategy and let the writing program build the flagship.",14,MIST,italic=True)
steps=[("1","Approve the tier→platform routing","Audible/Spotify/KDP-full get T7 (deep_book_6h); short tiers go to podcast / KDP short-read / Ximalaya. Config PR is ready, do-NOT-merge until you sign off."),
 ("2","Greenlight the 52k flagship","Route the word_target to the writing program: a real “Complete” (T7) book is the only thing that wins Audible + full KDP. This is the #1 content move."),
 ("3","Close the thin-render shortfall","Catalog books render ~5k vs target — a writing-program depth/fill task. Real length, ≥70% completion, never padding."),
 ("4","Gate CJK on a measured render","ja can proceed (anchored). zh + ko expansion ratios are unverified — measure on Pearl Star before shipping CJK durations."),
]
y=2.05
for n,h,b in steps:
    sp=rect(s,0.6,y,0.5,0.5,GOLD,shape=MSO_SHAPE.OVAL)
    tf2=sp.text_frame; tf2.margin_left=0;tf2.margin_right=0;tf2.margin_top=0;tf2.margin_bottom=0
    para(tf2,n,18,INK,bold=True,font=HDR,align=PP_ALIGN.CENTER,first=True,space_after=0)
    _,tf3=box(s,1.3,y-0.05,11.3,1.15)
    para(tf3,h,16,GOLD,bold=True,font=HDR,first=True,space_after=2)
    para(tf3,b,13,WHITE,space_after=0)
    y+=1.18
footer(s,9)

path=os.path.join(OUT,"DURATION_PER_PLATFORM_DECK.pptx")
prs.save(path)
print("saved",path,"slides",len(prs.slides._sldIdLst))
